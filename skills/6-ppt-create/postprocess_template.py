#!/usr/bin/env python3
"""postprocess_slides.py — fill in DATA BLOCK only; do not edit below the separator."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
try:
    from PIL import Image
except ImportError:
    Image = None

# ============================================================
# DATA BLOCK — fill in these values, nothing else
# ============================================================
PPTX_PATH = "output.pptx"          # <-- path to pandoc output

SECTION_TAGS = {
    # slide_index (0-based): ("tag_text", "color_name")
    # e.g.  1: ("z-space", "steelblue"),
}

LAYOUT_MAP = {
    "figure-dominant": [],   # slide indices with figure-dominant layout (annotated normalised in Phase 1)
    "vstack-3":        [],   # slide indices with vstack-3 layout
    "hstack-3":        [],   # slide indices with hstack-3 layout
    "table":           [],   # slide indices containing a stats table
}

FIGURES = {
    # slide_index: "/absolute/path/to/figure.png"
    # hstack-3 / vstack-3 entries: slide_index: ["/abs/fig_a.png", "/abs/fig_b.png", "/abs/fig_c.png"]
}

FIGURE_HEIGHT_OVERRIDE = {
    # optional, from explicit storyboard request
    # slide_index: 10/14,
}
# ============================================================

COLOR_MAP = {
    'steelblue':  (70, 130, 180),
    'darkorange': (255, 140,   0),
    'teal':       (  0, 128, 128),
    'purple':     (128,   0, 128),
    'crimson':    (220,  20,  60),
}

def get_aspect(path, fallback=1.6):
    if Image is None or path is None:
        return fallback
    try:
        w, h = Image.open(path).size
        return w / h
    except Exception:
        return fallback

def add_section_tag(slide, tag_text, color_name, slide_w):
    rgb = COLOR_MAP.get(color_name, (70, 130, 180))
    box = slide.shapes.add_textbox(Inches(0.08), Inches(0.05), int(slide_w * 0.70), Inches(0.36))
    tf  = box.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text           = tag_text
    run.font.size      = Pt(24)
    run.font.bold      = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*rgb)
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)

def clear_title(slide):
    title = slide.shapes.title
    if title is not None and title.has_text_frame:
        for para in title.text_frame.paragraphs:
            for run in para.runs:
                run.text = ""

def remove_empty_placeholders(slide):
    """Delete placeholder shapes whose text content is entirely empty."""
    sp_tree = slide.shapes._spTree
    for shape in list(slide.placeholders):
        if not shape.has_text_frame:
            continue
        all_text = ''.join(
            run.text for para in shape.text_frame.paragraphs for run in para.runs
        )
        if not all_text.strip():
            sp_tree.remove(shape._element)

def format_notes_one_sentence_per_line(slide):
    """Reformat speaker notes so each sentence occupies its own paragraph."""
    import re
    try:
        ntf = slide.notes_slide.notes_text_frame
    except Exception:
        return
    from pptx.oxml.ns import qn
    txBody = ntf._txBody
    paras  = txBody.findall(qn('a:p'))
    blocks = [''.join(r.text or '' for r in p.iter(qn('a:t'))) for p in paras]
    full = '\n\n'.join(b for b in blocks if b.strip())
    if not full.strip():
        return
    sentences = []
    for chunk in re.split(r'\n\n+', full):
        chunk = chunk.strip()
        if not chunk:
            continue
        sents = re.split(r'(?<=[.!?])\s+', chunk)
        sentences.extend(s.strip() for s in sents if s.strip())
    for p in paras[1:]:
        txBody.remove(p)
    first = paras[0]
    for r in list(first.findall(qn('a:r'))):
        first.remove(r)
    if sentences:
        ntf.paragraphs[0].text = sentences[0]
        for sent in sentences[1:]:
            ntf.add_paragraph().text = sent

def set_body_font(slide):
    title_id = slide.shapes.title.shape_id if slide.shapes.title else -1
    for shape in slide.shapes:
        if shape.shape_id == title_id:
            continue
        _set_text_shape_font(shape, 20)

def _set_text_shape_font(shape, font_pt):
    if not hasattr(shape, 'text_frame'):
        return
    for para in shape.text_frame.paragraphs:
        if para.runs:
            for run in para.runs:
                run.font.size = Pt(font_pt)
        else:
            para.font.size = Pt(font_pt)

def _boxes_overlap(a, b):
    ax1, ay1 = int(a.left), int(a.top)
    ax2, ay2 = ax1 + int(a.width), ay1 + int(a.height)
    bx1, by1 = int(b.left), int(b.top)
    bx2, by2 = bx1 + int(b.width), by1 + int(b.height)
    return (ax1 < bx2 and ax2 > bx1 and ay1 < by2 and ay2 > by1)

def _resolve_text_picture_overlap(slide, text_shape):
    if text_shape is None:
        return
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if not pics:
        return
    for font_pt in (20, 18, 16):
        _set_text_shape_font(text_shape, font_pt)
        if not any(_boxes_overlap(text_shape, pic) for pic in pics):
            return

def ensure_bullet_placeholder(slide, slide_w):
    """Ensure a top-half bullet band (1-3 bullets) below the section tag."""
    band_left = int(Inches(0.5))
    band_top = int(Inches(0.52))
    band_width = int(slide_w - 2 * Inches(0.5))
    band_height = int(Inches(1.6))

    bullet_shape = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            bullet_shape = ph
            break

    if bullet_shape is None:
        bullet_shape = slide.shapes.add_textbox(band_left, band_top, band_width, band_height)
    else:
        bullet_shape.left = band_left
        bullet_shape.top = band_top
        bullet_shape.width = band_width
        bullet_shape.height = band_height

    tf = bullet_shape.text_frame
    tf.word_wrap = True

    has_user_text = any(
        run.text.strip()
        for para in tf.paragraphs
        for run in para.runs
    )
    if not has_user_text:
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = "- [Bullet 1]"
        p1 = tf.add_paragraph()
        p1.text = "- [Bullet 2]"
        p2 = tf.add_paragraph()
        p2.text = "- [Bullet 3]"

    _set_text_shape_font(bullet_shape, 20)
    return bullet_shape

def _count_bullets(text_shape):
    if text_shape is None or not hasattr(text_shape, 'text_frame'):
        return 0
    count = 0
    for para in text_shape.text_frame.paragraphs:
        if para.text.strip():
            count += 1
    return count

def _default_figure_fraction_from_bullets(bullet_count):
    if bullet_count <= 1:
        return 11 / 14
    if bullet_count == 2:
        return 10 / 14
    return 9 / 14

def _figure_region_geometry(slide_idx, slide_h, text_shape):
    frac = FIGURE_HEIGHT_OVERRIDE.get(slide_idx)
    if frac is None:
        frac = _default_figure_fraction_from_bullets(_count_bullets(text_shape))
    frac = max(1e-3, min(float(frac), 1.0))
    fig_h = int(slide_h * frac)
    fig_top = int(slide_h) - fig_h
    min_top = int(text_shape.top + text_shape.height + Inches(0.08)) if text_shape is not None else 0
    fig_top = max(fig_top, min_top)
    fig_h = int(slide_h) - fig_top
    return fig_top, fig_h

def _fit_pictures_to_bottom_region(slide, slide_idx, slide_w, slide_h, text_shape):
    """Fit 1-4 pictures into the target bottom figure region while preserving arrangement."""
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if not (1 <= len(pics) <= 4):
        return
    fig_top, fig_h = _figure_region_geometry(slide_idx, slide_h, text_shape)
    region_left = int(Inches(0.15))
    region_w = int(slide_w - 2 * Inches(0.15))

    box_left = min(int(s.left) for s in pics)
    box_top = min(int(s.top) for s in pics)
    box_right = max(int(s.left + s.width) for s in pics)
    box_bottom = max(int(s.top + s.height) for s in pics)
    box_w = max(box_right - box_left, 1)
    box_h = max(box_bottom - box_top, 1)

    scale = min(region_w / box_w, fig_h / box_h)
    out_w = int(box_w * scale)
    out_h = int(box_h * scale)
    target_left = region_left + max((region_w - out_w) // 2, 0)
    target_top = fig_top + max((fig_h - out_h) // 2, 0)

    for shp in pics:
        rel_x = int(shp.left) - box_left
        rel_y = int(shp.top) - box_top
        shp.left = int(target_left + rel_x * scale)
        shp.top = int(target_top + rel_y * scale)
        shp.width = int(max(int(shp.width * scale), 1))
        shp.height = int(max(int(shp.height * scale), 1))

def style_table(slide, shape, slide_idx, slide_w, slide_h, has_bullets):
    tbl    = shape.table
    n_cols = len(tbl.columns)
    max_chars = [0] * n_cols
    for row in tbl.rows:
        for j, cell in enumerate(row.cells):
            max_chars[j] = max(max_chars[j], len(cell.text_frame.text))
    char_emu = int(Inches(0.12))
    min_emu  = 6 * char_emu
    widths   = [max(max_chars[j] * char_emu, min_emu) for j in range(n_cols)]
    total    = sum(widths)
    cap      = int(slide_w * 1.00)
    if total > cap:
        scale  = cap / total
        widths = [max(int(w * scale), min_emu) for w in widths]
        widths[-1] = max(cap - sum(widths[:-1]), min_emu)
        total = sum(widths)
    for j, col in enumerate(tbl.columns):
        col.width = widths[j]
    shape.width = total
    shape.left  = int((slide_w - total) / 2)

    text_shape = ensure_bullet_placeholder(slide, slide_w) if has_bullets else None
    top_min, avail_h = _figure_region_geometry(slide_idx, slide_h, text_shape)
    tbl_h   = sum(row.height for row in tbl.rows)
    shape.top = top_min + max((avail_h - tbl_h) // 2, 0)

    for row_idx, row in enumerate(tbl.rows):
        for cell in row.cells:
            tc = cell._tc
            for tag in ('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill',
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill'):
                for el in tc.findall('.//' + tag):
                    el.getparent().remove(el)
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(16)
                for run in para.runs:
                    run.font.size = Pt(16)
                    run.font.bold = (row_idx == 0)

def rebuild_vstack3(slide, fig_paths, tag_info, slide_w, slide_h, slide_idx):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)
    margin_side = Inches(0.15)
    max_w       = slide_w - 2 * margin_side
    if tag_info:
        add_section_tag(slide, tag_info[0], tag_info[1], slide_w)
    text_shape = ensure_bullet_placeholder(slide, slide_w)
    fig_top_min, avail_h = _figure_region_geometry(slide_idx, slide_h, text_shape)
    row_h = avail_h / 3
    for i, path in enumerate(fig_paths):
        aspect = get_aspect(path)
        img_w  = min(max_w, row_h * aspect)
        img_h  = img_w / aspect
        left   = margin_side + (max_w - img_w) / 2
        top    = fig_top_min + i * row_h + (row_h - img_h) / 2
        slide.shapes.add_picture(str(path), int(left), int(top), int(img_w), int(img_h))
    _resolve_text_picture_overlap(slide, text_shape)

def rebuild_hstack3(slide, fig_paths, tag_info, slide_w, slide_h, slide_idx):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)
    margin      = Inches(0.08)
    col_w       = (slide_w - 2 * margin) / 3
    if tag_info:
        add_section_tag(slide, tag_info[0], tag_info[1], slide_w)
    text_shape = ensure_bullet_placeholder(slide, slide_w)
    fig_top_min, avail_h = _figure_region_geometry(slide_idx, slide_h, text_shape)
    for i, path in enumerate(fig_paths):
        aspect      = get_aspect(path)
        img_h       = min(col_w / aspect, avail_h)
        img_w       = img_h * aspect
        vert_offset = (avail_h - img_h) / 2
        top         = fig_top_min + vert_offset
        left        = margin + i * col_w + (col_w - img_w) / 2
        slide.shapes.add_picture(str(path), int(left), int(top), int(img_w), int(img_h))
    _resolve_text_picture_overlap(slide, text_shape)

def reposition_figure_dominant(slide, fig_path, slide_w, slide_h, slide_idx):
    aspect      = get_aspect(fig_path)
    text_shape  = ensure_bullet_placeholder(slide, slide_w)
    fig_top_min, max_h = _figure_region_geometry(slide_idx, slide_h, text_shape)
    img_w       = int(slide_w - 2 * Inches(0.3))
    img_h       = min(int(img_w / aspect), max_h)
    if img_h == max_h:
        img_w = int(img_h * aspect)
    fig_top  = fig_top_min + max((max_h - img_h) // 2, 0)
    fig_left = int((slide_w - img_w) / 2)
    for shape in slide.shapes:
        if shape.shape_type == 13:
            shape.left   = fig_left
            shape.top    = fig_top
            shape.width  = img_w
            shape.height = img_h
            break
    _resolve_text_picture_overlap(slide, text_shape)

def main():
    prs     = Presentation(PPTX_PATH)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for idx, slide in enumerate(prs.slides):
        tag_info = SECTION_TAGS.get(idx)
        fig_path = FIGURES.get(idx)

        if idx in LAYOUT_MAP.get("vstack-3", []):
            paths = fig_path if isinstance(fig_path, list) else [fig_path]
            rebuild_vstack3(slide, paths, tag_info, slide_w, slide_h, idx)
            format_notes_one_sentence_per_line(slide)
            continue

        if idx in LAYOUT_MAP.get("hstack-3", []):
            paths = fig_path if isinstance(fig_path, list) else [fig_path]
            rebuild_hstack3(slide, paths, tag_info, slide_w, slide_h, idx)
            format_notes_one_sentence_per_line(slide)
            continue

        if tag_info:
            add_section_tag(slide, tag_info[0], tag_info[1], slide_w)
            clear_title(slide)

        remove_empty_placeholders(slide)
        set_body_font(slide)

        has_bullets = False
        bullet_shape = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                continue
            if ph.has_text_frame and any(
                run.text.strip()
                for para in ph.text_frame.paragraphs
                for run in para.runs
            ):
                has_bullets = True
                bullet_shape = ph
                break

        for shape in slide.shapes:
            if shape.shape_type == 19:
                style_table(slide, shape, idx, slide_w, slide_h, has_bullets)

        if idx in LAYOUT_MAP.get("figure-dominant", []) and fig_path:
            reposition_figure_dominant(slide, fig_path, slide_w, slide_h, idx)
        elif 1 <= len([s for s in slide.shapes if s.shape_type == 13]) <= 4:
            if bullet_shape is None and has_bullets:
                bullet_shape = ensure_bullet_placeholder(slide, slide_w)
            _fit_pictures_to_bottom_region(slide, idx, slide_w, slide_h, bullet_shape)
            _resolve_text_picture_overlap(slide, bullet_shape)

        format_notes_one_sentence_per_line(slide)

    prs.save(PPTX_PATH)
    print(f"Saved: {PPTX_PATH}")

if __name__ == "__main__":
    main()
